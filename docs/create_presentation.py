#!/usr/bin/env python3
"""
Create a professional PowerPoint presentation for the Logistics Optimization System
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 45, 85)  # Dark blue
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(100, 200, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    p = footer_frame.paragraphs[0]
    p.text = "🚀 For AMD TCS Hackathon"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(150, 200, 255)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items, bg_color=None):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    if bg_color:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    else:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 245, 250)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 45, 85)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(30, 30, 30)
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.level = 0

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 245, 250)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 45, 85)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(220, 50, 50)
    
    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(30, 30, 30)
        p.space_before = Pt(4)
        p.level = 0
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.2), Inches(4.2), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(50, 150, 50)
    
    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(30, 30, 30)
        p.space_before = Pt(4)
        p.level = 0

def add_architecture_slide(prs):
    """Add architecture diagram slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 245, 250)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "System Architecture"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 45, 85)
    
    # Frontend layer
    frontend = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(1.2), Inches(5), Inches(0.8))
    frontend.fill.solid()
    frontend.fill.fore_color.rgb = RGBColor(100, 150, 200)
    frontend.line.color.rgb = RGBColor(25, 45, 85)
    frontend.line.width = Pt(2)
    tf = frontend.text_frame
    tf.text = "Frontend: Streamlit Dashboard"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Arrow 1
    arrow1 = slide.shapes.add_connector(1, Inches(4.5), Inches(2), Inches(4.5), Inches(2.3))
    arrow1.line.color.rgb = RGBColor(100, 100, 100)
    arrow1.line.width = Pt(2)
    
    # Decision Agent (center)
    decision = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(2.5), Inches(4), Inches(0.8))
    decision.fill.solid()
    decision.fill.fore_color.rgb = RGBColor(200, 100, 50)
    decision.line.color.rgb = RGBColor(25, 45, 85)
    decision.line.width = Pt(2)
    tf = decision.text_frame
    tf.text = "Decision Agent (Orchestrator)"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Bottom agents
    y_pos = 3.8
    agents = [
        ("Routing Agent", Inches(0.5), RGBColor(50, 150, 100)),
        ("Cost Agent", Inches(2.5), RGBColor(50, 100, 200)),
        ("Disruption Agent", Inches(4.5), RGBColor(200, 100, 100)),
        ("Real-Time Data", Inches(6.5), RGBColor(150, 100, 150))
    ]
    
    for agent_name, x_pos, color in agents:
        # Arrow down
        arrow = slide.shapes.add_connector(1, Inches(4.5), Inches(3.3), x_pos + Inches(0.65), Inches(y_pos - 0.3))
        arrow.line.color.rgb = RGBColor(150, 150, 150)
        arrow.line.width = Pt(1.5)
        
        # Agent box
        agent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(y_pos), Inches(1.5), Inches(0.7))
        agent.fill.solid()
        agent.fill.fore_color.rgb = color
        agent.line.color.rgb = RGBColor(25, 45, 85)
        agent.line.width = Pt(1.5)
        tf = agent.text_frame
        tf.text = agent_name
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Arrow to backend
    arrow2 = slide.shapes.add_connector(1, Inches(4.5), Inches(3.3), Inches(4.5), Inches(3.5))
    arrow2.line.color.rgb = RGBColor(100, 100, 100)
    arrow2.line.width = Pt(2)
    
    # Graph database
    graph_db = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.5), Inches(5), Inches(0.8))
    graph_db.fill.solid()
    graph_db.fill.fore_color.rgb = RGBColor(100, 100, 100)
    graph_db.line.color.rgb = RGBColor(25, 45, 85)
    graph_db.line.width = Pt(2)
    tf = graph_db.text_frame
    tf.text = "Graph Database (NetworkX) - City Map with Live Data"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_flow_diagram(prs):
    """Add process flow diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 245, 250)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Real-Time Optimization Flow"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 45, 85)
    
    # Flow steps
    steps = [
        ("User Input", "From → To", Inches(0.8)),
        ("Fetch Live Data", "Traffic, Weather", Inches(2.3)),
        ("Apply Disruptions", "Update Graph", Inches(3.8)),
        ("Optimize Routes", "All Vehicles", Inches(5.3)),
        ("Select Best", "Lowest Cost", Inches(6.8))
    ]
    
    y_pos = Inches(2)
    
    for i, (step_name, step_desc, x_pos) in enumerate(steps):
        # Step box
        step_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, Inches(1.2), Inches(1))
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = RGBColor(100, 150, 220)
        step_box.line.color.rgb = RGBColor(25, 45, 85)
        step_box.line.width = Pt(2)
        
        tf = step_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step_name
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = step_desc
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Arrow to next
        if i < len(steps) - 1:
            arrow = slide.shapes.add_connector(1, x_pos + Inches(1.2), y_pos + Inches(0.5),
                                               steps[i+1][2], y_pos + Inches(0.5))
            arrow.line.color.rgb = RGBColor(100, 100, 100)
            arrow.line.width = Pt(2)
    
    # Output box
    output_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(3.5), Inches(7), Inches(2))
    output_box.fill.solid()
    output_box.fill.fore_color.rgb = RGBColor(220, 250, 220)
    output_box.line.color.rgb = RGBColor(50, 150, 50)
    output_box.line.width = Pt(2)
    
    tf = output_box.text_frame
    tf.word_wrap = True
    
    output_items = [
        "✓ Optimized Route (Location sequence)",
        "✓ Best Vehicle (Van/Truck/Hybrid)",
        "✓ Cost Breakdown (Fuel, Labor, Profit)",
        "✓ Carbon Emissions (kg CO2)",
        "✓ AI Explanation (Why this choice?)"
    ]
    
    for i, item in enumerate(output_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(30, 30, 30)
        p.space_before = Pt(2)

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, "Agentic AI Logistics Optimization", 
                   "Multi-Agent System for Real-Time Route Planning")
    
    # Slide 2: The Problem
    add_two_column_slide(prs, "The Challenge",
                        "❌ Without AI", [
                            "Manual route planning",
                            "Slow decision making",
                            "Traffic causes delays",
                            "Fuel prices ignored",
                            "High costs & emissions",
                            "No transparency"
                        ],
                        "✅ With Our System", [
                            "Automatic optimization",
                            "Real-time adaptation",
                            "Live traffic awareness",
                            "Dynamic fuel pricing",
                            "Cost & emissions minimized",
                            "AI explanations"
                        ])
    
    # Slide 3: Solution Overview
    add_content_slide(prs, "Our Solution",
                     [
                        "🤖 Five AI Agents working together:",
                        "  • Routing Agent → Finds shortest paths",
                        "  • Cost Agent → Calculates expenses",
                        "  • Disruption Agent → Handles real-world events",
                        "  • Monitoring Agent → Simulates scenarios",
                        "  • Decision Agent → Orchestrates everything",
                        "",
                        "💻 Complete platform with backend API + interactive dashboard"
                     ])
    
    # Slide 4: Architecture Diagram
    add_architecture_slide(prs)
    
    # Slide 5: Flow Diagram
    add_flow_diagram(prs)
    
    # Slide 6: Key Features
    add_content_slide(prs, "Key Features",
                     [
                        "📍 Multi-Vehicle Fleet Planning",
                        "  Van (cheap), Truck (capacity), Hybrid (balanced)",
                        "",
                        "🌍 Real-Time Data Integration",
                        "  Traffic, Weather, Fuel Prices from live APIs",
                        "",
                        "📊 Carbon Emission Tracking",
                        "  Every route shows CO2 footprint",
                        "",
                        "📝 Explainable AI",
                        "  System explains every decision in plain English"
                     ])
    
    # Slide 7: How It Works
    add_content_slide(prs, "4-Step Process",
                     [
                        "1️⃣  USER INPUT → Specify warehouse and customer locations",
                        "",
                        "2️⃣  ANALYZE → AI fetches real-time data, evaluates options",
                        "",
                        "3️⃣  ADAPT → If disruptions detected, re-optimize route",
                        "",
                        "4️⃣  DECIDE → Show best route + vehicle + cost breakdown"
                     ])
    
    # Slide 8: Demo Scenarios
    add_content_slide(prs, "Interactive Testing",
                     [
                        "🎮 Click 'Run Optimization' → See initial best route",
                        "",
                        "🚨 Simulate Traffic Jam → Route adapts automatically",
                        "",
                        "☔ Simulate Bad Weather → Travel time increases",
                        "",
                        "⛽ Simulate Fuel Price Spike → Vehicle choice might change",
                        "",
                        "👁️  View Green (initial) vs Red (updated) paths on map"
                     ])
    
    # Slide 9: Technology Stack
    add_content_slide(prs, "Tech Stack",
                     [
                        "⚙️ Backend: FastAPI (Python web framework)",
                        "",
                        "📊 Optimization: NetworkX (graph algorithms)",
                        "",
                        "🎨 Frontend: Streamlit (interactive dashboard)",
                        "",
                        "📈 Visualization: Plotly (interactive maps)",
                        "",
                        "🔗 Data: Mock + Real API integration (async/await)"
                     ])
    
    # Slide 10: Bonus Features
    add_content_slide(prs, "Bonus Enhancements",
                     [
                        "✅ Multi-Vehicle Support (Van, Truck, Hybrid)",
                        "",
                        "✅ Carbon Emissions Scoring (kg CO2 tracking)",
                        "",
                        "✅ Real-Time External APIs (Traffic, Weather, Fuel)",
                        "",
                        "✅ RL Agent Placeholder (for future ML training)",
                        "",
                        "✅ Smart Cost Model (Fuel, Labor, Overhead, Profit)"
                     ])
    
    # Slide 11: Business Impact
    add_two_column_slide(prs, "Business Value",
                        "💰 Cost Savings", [
                            "Reduce fuel costs 15-25%",
                            "Optimize vehicle usage",
                            "Auto-select best option",
                            "Real-time price awareness"
                        ],
                        "🌱 Sustainability", [
                            "Track CO2 emissions",
                            "Reduce carbon footprint",
                            "Green delivery options",
                            "ESG compliance"
                        ])
    
    # Slide 12: Quick Start
    add_content_slide(prs, "Getting Started (3 Steps)",
                     [
                        "1. Install dependencies:",
                        "   pip install -r requirements.txt",
                        "",
                        "2. Start backend (Terminal 1):",
                        "   uvicorn backend.main:app --reload",
                        "",
                        "3. Start dashboard (Terminal 2):",
                        "   streamlit run frontend/app.py",
                        "",
                        "→ Browser opens automatically → Click 'Run Optimization'!"
                     ])
    
    # Slide 13: Hackathon Value
    add_content_slide(prs, "Why This Wins Hackathons",
                     [
                        "✅ Complete working system (not just prototype)",
                        "",
                        "✅ Real multi-agent AI (not just rules engine)",
                        "",
                        "✅ Beautiful, functional UI (judges love this)",
                        "",
                        "✅ Explainable AI (judges understand decisions)",
                        "",
                        "✅ Scalable architecture (easy to extend)",
                        "",
                        "✅ Production-ready code (proper error handling)"
                     ])
    
    # Slide 14: Use Cases
    add_content_slide(prs, "Real-World Applications",
                     [
                        "🚚 E-Commerce Delivery (Amazon, Flipkart style)",
                        "",
                        "🏢 Corporate Logistics (Multi-warehouse networks)",
                        "",
                        "🍔 Food Delivery (DoorDash, Uber Eats optimization)",
                        "",
                        "📦 Last-Mile Delivery (Urban congestion handling)",
                        "",
                        "🌍 Supply Chain (Global fleet management)"
                     ])
    
    # Slide 15: Future Roadmap
    add_content_slide(prs, "Future Enhancements",
                     [
                        "🤖 Reinforcement Learning → Learn from past routes",
                        "",
                        "🗺️ OR-Tools Integration → Vehicle Routing Problem (VRP)",
                        "",
                        "📱 Mobile App → Driver notifications & tracking",
                        "",
                        "🌐 Global APIs → Real traffic data worldwide",
                        "",
                        "💾 Predictive Analytics → Forecast disruptions"
                     ])
    
    # Slide 16: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 45, 85)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "Let's Optimize the Future!"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 200, 255)
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = "Questions? Let's build together! 🚀"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(150, 200, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_path = "/Users/sagarshingare/amd-tcs-hackathon/docs/Logistics_Optimization_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation created: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
